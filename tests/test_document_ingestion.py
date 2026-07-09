"""Tests for multi-format document ingestion (docx/pptx/pdf/xls/images + OCR).

Fixtures are generated in-memory (python-docx / python-pptx / xlwt / Pillow) or
hand-crafted (a minimal valid text PDF), so the suite needs no checked-in binary
assets. The OCR paths are exercised with the Tesseract calls monkeypatched, so no
real ``tesseract`` binary is required in CI.
"""

from __future__ import annotations

import io
import zipfile

import pytest

from app.api import packs as packs_module
from app.api.packs import (
    UploadProcessingError,
    _extract_folder_zip,
    _is_folder_zip,
    _parse_upload_records,
    _records_from_docx,
    _records_from_pptx,
    _records_from_table,
)


# ---------------------------------------------------------------------------
# In-memory fixture builders
# ---------------------------------------------------------------------------
def _make_docx() -> bytes:
    import docx

    document = docx.Document()
    document.add_paragraph("Quarterly revenue increased by 12 percent.")
    document.add_paragraph("")  # empty paragraph should be dropped
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Account"
    table.cell(0, 1).text = "Balance"
    table.cell(1, 0).text = "Cash"
    table.cell(1, 1).text = "10000 JPY"
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _make_pptx() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Forecast summary for FY2026"
    slide.notes_slide.notes_text_frame.text = "Presenter note: emphasize FX risk"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _make_text_pdf(text: str) -> bytes:
    """Build a minimal, valid single-page PDF with real embedded text."""
    stream = f"BT /F1 24 Tf 72 700 Td ({text}) Tj ET".encode("latin-1")
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>"
        ),
        b"<< /Length %d >>\nstream\n%s\nendstream" % (len(stream), stream),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, body in enumerate(objs, start=1):
        offsets.append(len(out))
        out += b"%d 0 obj\n" % i + body + b"\nendobj\n"
    xref_pos = len(out)
    out += b"xref\n0 %d\n" % (len(objs) + 1)
    out += b"0000000000 65535 f \n"
    for off in offsets:
        out += b"%010d 00000 n \n" % off
    out += b"trailer\n<< /Size %d /Root 1 0 R >>\n" % (len(objs) + 1)
    out += b"startxref\n%d\n%%%%EOF" % xref_pos
    return bytes(out)


def _make_xls() -> bytes:
    import xlwt

    book = xlwt.Workbook()
    sheet = book.add_sheet("Sheet1")
    sheet.write(0, 0, "Title")
    sheet.write(0, 1, "ContentText")
    sheet.write(1, 0, "Invoice")
    sheet.write(1, 1, "Legacy xls total 250 USD")
    buffer = io.BytesIO()
    book.save(buffer)
    return buffer.getvalue()


def _make_png() -> bytes:
    from PIL import Image

    image = Image.new("RGB", (8, 8), color=(255, 255, 255))
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# Native text extraction (no OCR)
# ---------------------------------------------------------------------------
def test_docx_extracts_paragraphs_and_tables():
    records = _parse_upload_records("report.docx", _make_docx())
    joined = "\n".join(records)
    assert "Quarterly revenue increased by 12 percent." in joined
    # Table cells are extracted too.
    assert "Cash" in joined and "10000 JPY" in joined
    # Empty paragraph is not a record.
    assert all(r.strip() for r in records)


def test_pptx_extracts_shape_text_and_notes():
    records = _parse_upload_records("deck.pptx", _make_pptx())
    joined = "\n".join(records)
    assert "Forecast summary for FY2026" in joined
    assert "emphasize FX risk" in joined


def test_text_pdf_extraction():
    records = _parse_upload_records("invoice.pdf", _make_text_pdf("Invoice Total 12345 JPY"))
    joined = "\n".join(records)
    assert "Invoice Total 12345 JPY" in joined


def test_legacy_xls_extraction():
    records = _parse_upload_records("legacy.xls", _make_xls())
    joined = "\n".join(records)
    assert "Legacy xls total 250 USD" in joined
    # Header-labeled like the xlsx reader.
    assert "ContentText:" in joined or "Title:" in joined


# ---------------------------------------------------------------------------
# OCR paths (Tesseract mocked -- no binary required)
# ---------------------------------------------------------------------------
def test_pdf_falls_back_to_ocr_when_no_text(monkeypatch):
    calls = {"ocr": 0}

    def fake_text_pages(data):
        return [""]  # simulate a scanned/image PDF with no embedded text

    def fake_ocr_pdf(data):
        calls["ocr"] += 1
        return ["OCR extracted invoice text"]

    monkeypatch.setattr(packs_module, "_pdf_text_pages", fake_text_pages)
    monkeypatch.setattr(packs_module, "_ocr_pdf", fake_ocr_pdf)

    records = _parse_upload_records("scan.pdf", b"%PDF-1.4 fake")
    assert calls["ocr"] == 1
    assert records == ["OCR extracted invoice text"]


def test_image_upload_uses_ocr(monkeypatch):
    calls = {"ocr": 0}

    def fake_ocr_images(images):
        calls["ocr"] += 1
        assert len(images) == 1  # single image passed through
        return ["OCR text from image"]

    monkeypatch.setattr(packs_module, "_ocr_images", fake_ocr_images)

    records = _parse_upload_records("receipt.png", _make_png())
    assert calls["ocr"] == 1
    assert records == ["OCR text from image"]


def test_ocr_images_reports_missing_tesseract(monkeypatch):
    """A missing Tesseract binary surfaces as a clear UploadProcessingError."""
    import pytesseract

    def boom(*args, **kwargs):
        raise pytesseract.TesseractNotFoundError()

    monkeypatch.setattr(pytesseract, "image_to_string", boom)

    from PIL import Image

    with pytest.raises(UploadProcessingError) as exc:
        packs_module._ocr_images([Image.new("RGB", (4, 4))])
    assert exc.value.status_code == 500


# ---------------------------------------------------------------------------
# Unsupported types & folder-zip interaction
# ---------------------------------------------------------------------------
def test_unsupported_extension_still_415():
    with pytest.raises(UploadProcessingError) as exc:
        _parse_upload_records("malware.exe", b"MZ...")
    assert exc.value.status_code == 415


def test_single_docx_not_detected_as_folder_zip():
    """A single .docx (which is itself a ZIP) must not look like a folder ZIP."""
    assert _is_folder_zip(_make_docx()) is False
    assert _is_folder_zip(_make_pptx()) is False


def test_folder_zip_with_docx_member_is_ingested():
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("folder/report.docx", _make_docx())
        archive.writestr("folder/notes.txt", b"Plain text note about revenue.")
    payload = buffer.getvalue()

    assert _is_folder_zip(payload) is True
    records, sources, warnings, unsupported = _extract_folder_zip(payload)
    joined = "\n".join(records)
    assert "Quarterly revenue increased by 12 percent." in joined
    assert "Plain text note about revenue." in joined
    assert "report.docx" in sources
    assert unsupported == []


# ---------------------------------------------------------------------------
# Table -> "header: value" record quality (retrieval-friendly)
# ---------------------------------------------------------------------------
def _make_docx_keyvalue_table() -> bytes:
    """A docx with a paragraph + a 2-column key/value spec table (Field|Value)."""
    import docx

    document = docx.Document()
    document.add_paragraph("Program brief for the new sportswear line.")
    rows = [
        ("Field", "Value"),
        ("Budget cap", "USD 420,000"),
        ("Pilot market", "Japan and Singapore"),
    ]
    table = document.add_table(rows=len(rows), cols=2)
    for r, (k, v) in enumerate(rows):
        table.cell(r, 0).text = k
        table.cell(r, 1).text = v
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _make_pptx_with_table() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(4), Inches(0.8))
    box.text_frame.text = "Launch metrics"
    rows, cols = 3, 2
    graphic = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.5), Inches(6), Inches(2)
    )
    table = graphic.table
    cells = [
        ("Field", "Value"),
        ("Budget cap", "USD 420,000"),
        ("Pilot market", "Japan and Singapore"),
    ]
    for r, (k, v) in enumerate(cells):
        table.cell(r, 0).text = k
        table.cell(r, 1).text = v
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_docx_keyvalue_table_becomes_header_value_records():
    records = _records_from_docx(_make_docx_keyvalue_table())
    # Paragraph is preserved (regression).
    assert "Program brief for the new sportswear line." in records
    # Key/value rows become retrieval-friendly "header: value" records.
    assert "Budget cap: USD 420,000" in records
    assert "Pilot market: Japan and Singapore" in records
    # The generic "Field | Value" header row is not emitted as a record.
    assert "Field: Value" not in records


def test_pptx_table_becomes_header_value_records():
    records = _records_from_pptx(_make_pptx_with_table())
    joined = "\n".join(records)
    # Non-table shape text is still captured.
    assert "Launch metrics" in joined
    # Table rows become "header: value" records.
    assert "Budget cap: USD 420,000" in records
    assert "Pilot market: Japan and Singapore" in records


def test_records_from_table_wide_uses_column_headers():
    rows = [
        ["Region", "Q1", "Q2"],
        ["APAC", "120", "150"],
    ]
    records = _records_from_table(rows)
    assert records == ["Region: APAC\nQ1: 120\nQ2: 150"]


def test_records_from_table_column_without_header_keeps_bare_value():
    rows = [
        ["Region", "Q1", ""],
        ["APAC", "120", "note"],
    ]
    records = _records_from_table(rows)
    assert records == ["Region: APAC\nQ1: 120\nnote"]

