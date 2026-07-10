"""Memory Forge + Memory Lens endpoints for Safe Memory Packs."""

from __future__ import annotations

import asyncio
import io
import json
import logging
import re
import secrets
import uuid
import zipfile
from collections import Counter
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional

from fastapi import APIRouter, BackgroundTasks, Depends, File, Form, HTTPException, UploadFile
from pydantic import ValidationError

from app.config import settings
from app.core import export_links, jobs_store, oss_storage, pack_import, pack_io, storage, upload_links, url_normalize
from app.core.audit import audit_export_result, build_audit_report
from app.core.auth import UploadAuthContext, require_upload_or_token
from app.core.chunker import chunk_text
from app.core.ledger import append_ledger_block, verify_ledger_chain
from app.core.pack_io import UnsafePathError
from app.core.policy import (
    can_export_entry,
    can_send_entry_to_llm,
    can_use_entry_for_query,
    classification_to_policy,
    redact_text_if_needed,
)
from app.core.qwen_client import qwen_client
from app.core.search import hybrid_search
from app.core.translation import normalize_accounting_batch
from app.models.job_schema import (
    JobRecord,
    JobResponse,
    JobStatus,
    RetentionMode,
    job_to_response,
)
from app.models.pack_schema import (
    AppendRequest,
    AppendResponse,
    BuildFromUploadRefRequest,
    BuildFromUrlRequest,
    BuildPackRequest,
    BuildPackResponse,
    BuildRefAcceptedResponse,
    Classification,
    Entry,
    ExportRequest,
    ExportResponse,
    ImportByRefRequest,
    ImportByRefResponse,
    ImportFromUploadRefRequest,
    Manifest,
    Metadata,
    Provenance,
    QueryHit,
    QueryRequest,
    QueryResponse,
    QueryByUploadRequest,
    QueryByUploadResponse,
    SafeMemoryPack,
    UploadBuildResponse,
    UploadedPackHit,
    VerifyRequest,
    VerifyResponse,
    get_retrieval_text,
)
from app.models.upload_schema import UploadRecord, UploadStatus
from app.core.pack_import import ImportRefError

router = APIRouter(prefix="/api/packs", tags=["packs"])

logger = logging.getLogger("safe_memory.packs")

# Maximum accepted upload size for /build-from-upload (bytes).
# Legacy default kept for reference; the effective limit is configurable via
# SAFE_MEMORY_MAX_UPLOAD_MB (see _max_upload_bytes).
MAX_UPLOAD_BYTES = 5 * 1024 * 1024  # 5 MB


def _max_upload_bytes() -> int:
    """Effective max upload size in bytes from configuration."""
    return max(1, int(settings.safe_memory_max_upload_mb)) * 1024 * 1024


def _max_import_bytes() -> int:
    """Effective max remote-import size in bytes from configuration."""
    return max(1, int(settings.safe_memory_max_import_mb)) * 1024 * 1024


def _absolute_url(rel_path: str) -> str:
    """Build an absolute HTTPS URL from the public base, else return relative."""
    base = (settings.safe_memory_public_base_url or "").strip().rstrip("/")
    if not rel_path.startswith("/"):
        rel_path = "/" + rel_path
    return f"{base}{rel_path}" if base else rel_path

_CJK_RE = re.compile(r"[\u3040-\u30ff\u3400-\u4dbf\u4e00-\u9fff]")

_STOPWORDS = {
    "the", "and", "for", "are", "but", "not", "you", "all", "any", "can",
    "her", "was", "one", "our", "out", "has", "had", "his", "she", "him",
    "this", "that", "with", "from", "they", "them", "then", "than", "have",
    "will", "your", "into", "some", "more", "what", "when", "which", "were",
    "there", "their", "about", "would", "these", "could", "other",
}
_WORD_RE = re.compile(r"[A-Za-z][A-Za-z0-9\-']+")


def extract_keywords(text: str, limit: int = 8) -> List[str]:
    """Extract simple frequency-based keywords from text."""
    words = [w.lower() for w in _WORD_RE.findall(text or "")]
    words = [w for w in words if len(w) > 2 and w not in _STOPWORDS]
    if not words:
        return []
    counts = Counter(words)
    return [w for w, _ in counts.most_common(limit)]


def _resolve_pack_path(
    agent_id: str,
    pack_path: Optional[str],
    pack_id: Optional[str],
):
    """Resolve a pack path from either an explicit path or a pack id."""
    if pack_path:
        return pack_io.ensure_safe_path(pack_path)
    if pack_id:
        found = pack_io.find_pack_by_id(agent_id, pack_id)
        if found is None:
            raise HTTPException(status_code=404, detail="Pack not found for pack_id.")
        return found
    raise HTTPException(status_code=400, detail="Provide pack_path or pack_id.")


def build_pack_from_entries(
    agent_id: str,
    pack_id: str,
    title: str,
    entries: List[Dict[str, Any]],
    default_classification: Classification = Classification.INTERNAL,
    method: str = "build",
    target_path: Optional[Any] = None,
):
    """Build and persist a Safe Memory Pack from raw entry specs.

    This is the reusable core used by both the ``/build`` endpoint and the
    Excel importer. Each item in ``entries`` is a dict that may contain:
    ``text``, ``original_text``, ``canonical_text``, ``source_language``,
    ``canonical_language``, ``translation_note``, ``classification``,
    ``source`` and ``tags``.

    Retrieval text (canonical English when present) drives embedding, keyword
    and classification generation. The original text is preserved for
    provenance and audit.

    Returns ``(pack, saved_path, audit_path)``.
    """
    # Determine retrieval text per spec; skip specs with no usable content.
    prepared: List[Dict[str, Any]] = []
    for spec in entries:
        canonical = (spec.get("canonical_text") or "").strip()
        base_text = (spec.get("text") or "").strip()
        retrieval_text = canonical or base_text
        if not retrieval_text:
            continue
        spec = dict(spec)
        spec["_retrieval_text"] = retrieval_text
        spec["_text"] = base_text or canonical
        prepared.append(spec)

    if not prepared:
        raise ValueError("No usable entries to build a pack.")

    embeddings = qwen_client.embed_texts([s["_retrieval_text"] for s in prepared])
    embedding_dim = len(embeddings[0]) if embeddings else 0

    manifest = Manifest(
        pack_id=pack_id,
        agent_id=agent_id,
        title=title,
        default_classification=default_classification,
        embedding_model=settings.qwen_embedding_model,
        chat_model=settings.qwen_chat_model,
        embedding_dim=embedding_dim,
    )
    pack = SafeMemoryPack(manifest=manifest, entries=[], ledger=[])

    for index, (spec, embedding) in enumerate(zip(prepared, embeddings)):
        retrieval_text = spec["_retrieval_text"]

        classification = spec.get("classification")
        if not isinstance(classification, Classification):
            classification = qwen_client.classify_text(
                retrieval_text, default=default_classification
            )

        entry = Entry(
            id=str(uuid.uuid4()),
            text=spec["_text"],
            embedding=embedding,
            keywords=extract_keywords(retrieval_text),
            classification=classification,
            policy=classification_to_policy(classification),
            metadata=Metadata(
                chunk_index=index,
                char_count=len(retrieval_text),
                tags=list(spec.get("tags") or []),
            ),
            provenance=Provenance(
                source=spec.get("source", "memory_forge"),
                origin_pack_id=pack_id,
                method=method,
            ),
            original_text=spec.get("original_text"),
            canonical_text=spec.get("canonical_text"),
            source_language=spec.get("source_language"),
            canonical_language=spec.get("canonical_language"),
            translation_note=spec.get("translation_note"),
        )
        pack.entries.append(entry)
        append_ledger_block(pack.ledger, entry, action=method)

    target = target_path or pack_io.pack_target_path(
        agent_id, pack_id, default_classification
    )
    saved_path = pack_io.save_pack(pack, target)

    report = build_audit_report(pack)
    audit_path = pack_io.save_audit_report(agent_id, pack_id, report)

    return pack, saved_path, audit_path


@router.post(
    "/build",
    response_model=BuildPackResponse,
    operation_id="buildMemoryPack",
    summary="Build a Safe Memory Pack from raw text",
    description=(
        "Create a new portable, policy-aware Safe Memory Pack from temporary "
        "source text. The text is chunked, embedded, classified, and sealed into "
        "an append-only pack file. Returns the new pack_id and entry count."
    ),
)
def build_pack(req: BuildPackRequest) -> BuildPackResponse:
    """Build a Safe Memory Pack from temporary source text."""
    if not req.source_text.strip():
        raise HTTPException(status_code=400, detail="source_text must not be empty.")

    chunks = chunk_text(req.source_text)
    if not chunks:
        raise HTTPException(status_code=400, detail="No content to chunk.")

    entry_specs = [
        {"text": chunk, "source": "memory_forge"} for chunk in chunks
    ]

    try:
        pack, saved_path, audit_path = build_pack_from_entries(
            agent_id=req.agent_id,
            pack_id=req.pack_id,
            title=req.title,
            entries=entry_specs,
            default_classification=req.default_classification,
            method="build",
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc))

    # delete_source_after_build: the source text was only held in memory and is
    # never persisted, so honoring the flag means we simply drop the reference.
    if req.delete_source_after_build:
        req.source_text = ""

    return BuildPackResponse(
        pack_id=req.pack_id,
        pack_path=pack_io.relpath_from_root(saved_path),
        entry_count=len(pack.entries),
        audit_path=pack_io.relpath_from_root(audit_path),
        classification=req.default_classification,
    )


@router.post(
    "/query",
    response_model=QueryResponse,
    operation_id="queryMemoryPack",
    summary="Query a Safe Memory Pack",
    description=(
        "Answer a natural-language question over a Safe Memory Pack using "
        "hybrid (embedding + keyword) search plus Qwen reasoning. SECRET is "
        "never returned or sent to the LLM; private text only with "
        "include_private=true. Use the agent_id and pack_id from "
        "getUploadLinkResult (or a known pack_id)."
    ),
)
def query_pack(req: QueryRequest) -> QueryResponse:
    """Query a Safe Memory Pack using hybrid search and Qwen reasoning."""
    if not req.query.strip():
        raise HTTPException(status_code=400, detail="query must not be empty.")

    # Opportunistic hygiene: drop any expired one-time-import packs before we
    # resolve this one, so an expired ephemeral pack is gone (-> 404) and never
    # answered from. Best-effort; never blocks the query.
    jobs_store.sweep_expired_quietly()

    try:
        path = _resolve_pack_path(req.agent_id, req.pack_path, req.pack_id)
        pack = pack_io.load_pack(path)
    except UnsafePathError:
        raise HTTPException(status_code=400, detail="Unsafe pack path rejected.")
    except FileNotFoundError:
        raise HTTPException(status_code=404, detail="Pack not found.")

    pack_id = pack.manifest.pack_id
    allowed = req.allowed_classifications
    usable = [e for e in pack.entries if can_use_entry_for_query(e, allowed)]

    warnings: List[str] = []

    if not usable:
        logger.info(
            "queryMemoryPack agent_id=%s pack_id=%s top_k=%s hits=0 "
            "secret_excluded=0 usable_entries=0 confidence=0.0 fallback=true",
            req.agent_id,
            pack_id,
            req.top_k,
        )
        return QueryResponse(
            answer="No memory entries are permitted for this query under the "
            "current policy and allowed classifications.",
            used_memory_ids=[],
            hits=[],
            classifications=[],
            fallback_used=True,
            pack_id=pack_id,
            confidence=0.0,
            warnings=["No usable entries for this query."],
        )

    query_embedding = qwen_client.embed_text(req.query)
    ranked = hybrid_search(req.query, query_embedding, usable, top_k=req.top_k)

    hits: List[QueryHit] = []
    llm_entries = []
    secret_excluded = 0
    for entry, score in ranked:
        # Never expose SECRET content in a public-facing response.
        if entry.classification == Classification.SECRET:
            secret_excluded += 1
            continue

        retrieval_text = get_retrieval_text(entry)
        hit = QueryHit(
            entry_id=entry.id,
            text=retrieval_text,
            score=score,
            classification=entry.classification,
        )
        # Original (often private, non-English) text is only surfaced when the
        # authorized caller explicitly opts in.
        if req.include_private and entry.original_text:
            hit.original_text = entry.original_text
        hits.append(hit)

        # SECRET entries never leave the machine via an LLM call.
        if can_send_entry_to_llm(entry):
            llm_entries.append({"id": entry.id, "text": retrieval_text})

    if secret_excluded:
        warnings.append(
            f"{secret_excluded} SECRET entr"
            f"{'y was' if secret_excluded == 1 else 'ies were'} excluded from "
            "this response and never sent to the LLM."
        )

    result = qwen_client.answer_with_context(req.query, llm_entries)

    confidence = round(max((h.score for h in hits), default=0.0), 4)

    logger.info(
        "queryMemoryPack agent_id=%s pack_id=%s top_k=%s hits=%d "
        "secret_excluded=%d usable_entries=%d confidence=%s fallback=%s",
        req.agent_id,
        pack_id,
        req.top_k,
        len(hits),
        secret_excluded,
        len(usable),
        confidence,
        bool(result["fallback_used"]),
    )

    return QueryResponse(
        answer=str(result["answer"]),
        used_memory_ids=[h.entry_id for h in hits],
        hits=hits,
        classifications=[h.classification for h in hits],
        fallback_used=bool(result["fallback_used"]),
        pack_id=pack_id,
        confidence=confidence,
        warnings=warnings,
    )


@router.post(
    "/query-by-upload",
    response_model=QueryByUploadResponse,
    operation_id="queryUploadedMemory",
    summary="Search all packs uploaded via one import link",
    description=(
        "Search every pack a user uploaded through one import link. Pass the "
        "claim_id from createUploadLink(mode=import); you never need each pack's "
        "imp- agent_id/pack_id. Returns one merged answer plus used_packs, "
        "classifications, confidence, fallback. SECRET is never returned."
    ),
)
def query_uploaded_memory(req: QueryByUploadRequest) -> QueryByUploadResponse:
    """Cross-search all ephemeral packs bound to one import link's claim_id.

    Robustness helper for GPT Actions: the assistant keeps only the single
    ``claim_id`` from createUploadLink(mode=import) and never juggles the
    per-pack ``imp-`` agent_id/pack_id values (a prior source of empty 404
    queries). The server resolves the claim's private namespace, hybrid-searches
    each imported pack, merges hits by score, and synthesizes one answer. Tenant
    isolation: only this claim's packs are reachable.
    """
    if not req.query.strip():
        raise HTTPException(status_code=400, detail="query must not be empty.")

    # Opportunistic hygiene: drop expired ephemeral packs before resolving, so an
    # expired link answers nothing (mirrors query_pack). Best-effort.
    jobs_store.sweep_expired_quietly()

    claim = upload_links.load_claim(req.claim_id)
    if claim is None or claim.is_expired():
        raise HTTPException(status_code=404, detail="Unknown or expired claim_id.")

    mode = (getattr(claim, "mode", "build") or "build")
    if mode != "import":
        raise HTTPException(
            status_code=400,
            detail="This link is not an import link; nothing to query.",
        )

    # Only this claim's private, unguessable namespace is reachable. Other links
    # (and the shared vault) are never touched.
    agent_id = (claim.import_agent_id or "").strip()
    pack_ids: List[str] = []
    for item in (claim.imported or []):
        pid = str(item.get("pack_id") or "").strip()
        if pid and pid not in pack_ids:
            pack_ids.append(pid)

    query_embedding = qwen_client.embed_text(req.query)

    # NOTE on include_private: in queryMemoryPack, include_private only ever gates
    # surfacing a per-hit original_text field. This endpoint deliberately returns
    # aggregated used_packs + one synthesized answer and never exposes per-entry or
    # original_text (spec: original_text/SECRET non-exposed), so include_private has
    # no exposure effect here. The field is accepted for request-shape parity.
    _ = req.include_private

    # Rank each imported pack independently, then merge by score.
    scored: List[tuple] = []  # (entry, score, pack_id)
    searched_packs: List[str] = []
    for pack_id in pack_ids:
        if not agent_id:
            continue
        found = pack_io.find_pack_by_id(agent_id, pack_id)
        if found is None:
            continue
        try:
            pack = pack_io.load_pack(found)
        except Exception:  # noqa: BLE001 - skip unreadable packs, never fail the query.
            continue
        searched_packs.append(pack_id)
        usable = [e for e in pack.entries if can_use_entry_for_query(e, None)]
        if not usable:
            continue
        ranked = hybrid_search(req.query, query_embedding, usable, top_k=req.top_k)
        for entry, score in ranked:
            scored.append((entry, score, pack_id))

    # Global top_k across all packs.
    scored.sort(key=lambda t: t[1], reverse=True)
    top = scored[: req.top_k]

    llm_entries: List[Dict[str, str]] = []
    classifications: List[Classification] = []
    retained_scores: List[float] = []
    pack_hit_counts: Dict[str, int] = {}
    secret_excluded = 0
    for entry, score, pack_id in top:
        # Never expose or send SECRET content (parity with queryMemoryPack).
        if entry.classification == Classification.SECRET:
            secret_excluded += 1
            continue
        classifications.append(entry.classification)
        retained_scores.append(score)
        pack_hit_counts[pack_id] = pack_hit_counts.get(pack_id, 0) + 1
        if can_send_entry_to_llm(entry):
            llm_entries.append(
                {"id": entry.id, "text": get_retrieval_text(entry)}
            )

    result = qwen_client.answer_with_context(req.query, llm_entries)
    confidence = round(max(retained_scores, default=0.0), 4)

    # Per-pack structured log (never logs the API key, query body, or original
    # text). One line per searched pack for auditability.
    for pack_id in searched_packs:
        logger.info(
            "queryUploadedMemory claim_pack agent_id=%s pack_id=%s top_k=%s hits=%d",
            agent_id,
            pack_id,
            req.top_k,
            pack_hit_counts.get(pack_id, 0),
        )

    used_packs = [
        UploadedPackHit(
            agent_id=agent_id, pack_id=pack_id, hits=pack_hit_counts[pack_id]
        )
        for pack_id in searched_packs
        if pack_hit_counts.get(pack_id, 0) > 0
    ]

    return QueryByUploadResponse(
        answer=str(result["answer"]),
        used_packs=used_packs,
        classifications=classifications,
        confidence=confidence,
        fallback=bool(result["fallback_used"]),
    )


@router.post("/append", response_model=AppendResponse)
def append_entry(req: AppendRequest) -> AppendResponse:
    """Append a new memory entry to an existing pack."""
    if not req.text.strip():
        raise HTTPException(status_code=400, detail="text must not be empty.")

    try:
        path = pack_io.ensure_safe_path(req.pack_path)
        pack = pack_io.load_pack(path)
    except UnsafePathError:
        raise HTTPException(status_code=400, detail="Unsafe pack path rejected.")
    except FileNotFoundError:
        raise HTTPException(status_code=404, detail="Pack not found.")

    classification = qwen_client.classify_text(
        req.text,
        default=req.suggested_classification or Classification.INTERNAL,
    )
    embedding = qwen_client.embed_text(req.text)

    entry = Entry(
        id=str(uuid.uuid4()),
        text=req.text,
        embedding=embedding,
        keywords=extract_keywords(req.text),
        classification=classification,
        policy=classification_to_policy(classification),
        metadata=Metadata(
            chunk_index=len(pack.entries), char_count=len(req.text)
        ),
        provenance=Provenance(
            source=req.source,
            origin_pack_id=pack.manifest.pack_id,
            method="append",
        ),
    )
    pack.entries.append(entry)
    block = append_ledger_block(pack.ledger, entry, action="append")

    pack.manifest.version = _bump_patch(pack.manifest.version)
    pack_io.save_pack(pack, path)

    return AppendResponse(
        entry_id=entry.id,
        ledger_block_id=block.id,
        version=pack.manifest.version,
        classification=classification,
    )


@router.post(
    "/export",
    response_model=ExportResponse,
    operation_id="exportMemoryPack",
    summary="Export a safe, shareable copy of a pack",
    description=(
        "Produce a filtered, policy-aware copy of a pack for sharing. "
        "CONFIDENTIAL and SECRET entries are excluded unless explicitly allowed, "
        "sensitive text can be redacted, and private source text can be removed. "
        "Identify the pack with pack_id (preferred) or pack_path."
    ),
)
def export_pack(req: ExportRequest) -> ExportResponse:
    """Export a filtered, policy-aware copy of a pack."""
    try:
        path = _resolve_pack_path(req.agent_id, req.pack_path, req.pack_id)
        pack = pack_io.load_pack(path)
    except UnsafePathError:
        raise HTTPException(status_code=400, detail="Unsafe pack path rejected.")
    except FileNotFoundError:
        raise HTTPException(status_code=404, detail="Pack not found.")

    allowed = req.allowed_classifications
    included_entries: List[Entry] = []
    excluded_entries: List[Entry] = []

    for entry in pack.entries:
        if not can_export_entry(entry, allowed):
            excluded_entries.append(entry)
            continue

        exported = entry.model_copy(deep=True)
        if req.redact_sensitive_text:
            exported.text = redact_text_if_needed(exported.text, exported)
            if exported.canonical_text:
                exported.canonical_text = redact_text_if_needed(
                    exported.canonical_text, exported
                )
            if exported.original_text:
                exported.original_text = redact_text_if_needed(
                    exported.original_text, exported
                )
        if req.remove_sources:
            # Drop the private Japanese source and provenance for public export.
            exported.original_text = None
            exported.translation_note = None
            exported.provenance.source = "removed"
            exported.provenance.origin_pack_id = None
        included_entries.append(exported)

    # Rebuild a fresh ledger over the exported entries.
    export_manifest = pack.manifest.model_copy(deep=True)
    export_manifest.pack_id = f"{pack.manifest.pack_id}-export"
    export_manifest.title = f"{pack.manifest.title} (export)"
    export_manifest.updated_at = datetime.now(timezone.utc).isoformat()

    export_pack_obj = SafeMemoryPack(
        manifest=export_manifest, entries=included_entries, ledger=[]
    )
    for entry in included_entries:
        append_ledger_block(export_pack_obj.ledger, entry, action="export")

    target = pack_io.export_target_path(req.agent_id, req.export_name)
    saved = pack_io.save_pack(export_pack_obj, target)
    saved_rel = pack_io.relpath_from_root(saved)

    report = audit_export_result(
        included_entries,
        excluded_entries,
        allowed,
        req.redact_sensitive_text,
        req.remove_sources,
    )

    # Tokenized, API-key-free download link so another agent can fetch and
    # re-import this pack over plain HTTPS.
    token = export_links.create_export_link(req.agent_id, saved_rel)
    download_url = _absolute_url(f"/api/packs/dl/{token}")

    return ExportResponse(
        export_path=saved_rel,
        included_count=len(included_entries),
        excluded_count=len(excluded_entries),
        warnings=report["warnings"],
        download_url=download_url,
    )


@router.get(
    "/dl/{token}",
    include_in_schema=False,
)
def download_export(token: str):
    """Serve a pack by its download token (no API key required).

    Streams the local pack when it still exists; otherwise, if the token carries
    an ``oss_object_key`` and OSS is enabled, 307-redirects to a freshly signed
    OSS URL (the signature travels only in the ``Location`` header, so it is
    never corrupted by a text-rendering layer). Returns 404 when neither the
    local pack nor an OSS object is available.
    """
    record = export_links.resolve_export_link_record(token)
    if not record:
        raise HTTPException(status_code=404, detail="Download link not found.")

    rel_path = record.get("rel_path")
    if rel_path:
        try:
            path = pack_io.ensure_safe_path(rel_path)
        except UnsafePathError:
            raise HTTPException(status_code=400, detail="Unsafe pack path rejected.")
        if path.exists():
            from fastapi.responses import FileResponse

            return FileResponse(
                path=str(path),
                media_type="application/json",
                filename=path.name,
            )

    oss_object_key = record.get("oss_object_key")
    if oss_object_key:
        try:
            if oss_storage.is_enabled():
                signed = oss_storage.generate_signed_download_url(oss_object_key)
                from fastapi.responses import RedirectResponse

                return RedirectResponse(url=signed, status_code=307)
        except Exception:  # pragma: no cover - defensive; fall through to 404
            pass

    raise HTTPException(
        status_code=404, detail="Exported pack is no longer available."
    )


@router.post(
    "/verify",
    response_model=VerifyResponse,
    operation_id="verifyMemoryPack",
    summary="Verify a pack's integrity",
    description=(
        "Verify a pack's manifest and its append-only ledger sha256 hash chain "
        "so tampering is detectable. Identify the pack with pack_id + agent_id "
        "(preferred) or pack_path."
    ),
)
def verify_pack(req: VerifyRequest) -> VerifyResponse:
    """Verify a pack's manifest and ledger hash chain."""
    try:
        path = _resolve_pack_path(req.agent_id, req.pack_path, req.pack_id)
        pack = pack_io.load_pack(path)
    except UnsafePathError:
        raise HTTPException(status_code=400, detail="Unsafe pack path rejected.")
    except FileNotFoundError:
        raise HTTPException(status_code=404, detail="Pack not found.")

    manifest_present = bool(pack.manifest and pack.manifest.pack_id)
    valid_chain, warnings = verify_ledger_chain(pack.ledger, pack.entries)

    if not manifest_present:
        warnings.append("Manifest is missing or incomplete.")

    return VerifyResponse(
        valid_hash_chain=valid_chain,
        manifest_present=manifest_present,
        entry_count=len(pack.entries),
        ledger_count=len(pack.ledger),
        warnings=warnings,
    )


def _import_pack_from_bytes(
    raw: bytes, agent_id: str, pack_id: Optional[str], json_error_detail: str
) -> ImportByRefResponse:
    """Validate, verify, and import a ``.smp.json`` pack from raw bytes.

    Shared by :func:`import_pack_by_ref` (URL fetch) and
    :func:`import_pack_from_upload_ref` (staged upload). Re-homes the pack under
    ``agent_id`` so it becomes queryable by ``pack_id`` in that agent's vault.
    """
    try:
        data = json.loads(raw.decode("utf-8"))
    except (UnicodeDecodeError, json.JSONDecodeError):
        raise HTTPException(status_code=400, detail=json_error_detail)

    try:
        pack = SafeMemoryPack.model_validate(data)
    except ValidationError:
        raise HTTPException(
            status_code=400, detail="Content is not a valid Safe Memory Pack."
        )

    valid, ledger_warnings = verify_ledger_chain(pack.ledger, pack.entries)
    warnings = list(ledger_warnings)
    if not valid:
        if settings.safe_memory_import_require_valid_ledger:
            raise HTTPException(
                status_code=422,
                detail="Imported pack failed ledger verification (strict mode).",
            )
        warnings.append("Ledger verification failed; imported with verified=false.")

    final_pack_id = (
        (pack_id or "").strip()
        or (pack.manifest.pack_id or "").strip()
        or f"imported-{uuid.uuid4().hex[:8]}"
    )

    # Re-home the pack under the importing agent.
    pack.manifest.agent_id = agent_id
    pack.manifest.pack_id = final_pack_id

    try:
        target = pack_io.pack_target_path(
            agent_id, final_pack_id, pack.manifest.default_classification
        )
        pack_io.save_pack(pack, target)
    except UnsafePathError:
        raise HTTPException(status_code=400, detail="Unsafe pack path rejected.")

    counts = dict(Counter(e.classification.value for e in pack.entries))
    return ImportByRefResponse(
        pack_id=final_pack_id,
        entry_count=len(pack.entries),
        classification_summary=counts,
        verified=valid,
        warnings=warnings,
    )


@router.post(
    "/import-by-ref",
    response_model=ImportByRefResponse,
    operation_id="importPackByRef",
    summary="Import a Safe Memory Pack from a URL",
    description=(
        "Import a Safe Memory Pack from a public HTTPS URL (a .smp.json file, "
        "e.g. a server or Drive/OneDrive share link) into an agent's vault, so "
        "packs are shared by link. HTTPS-only with SSRF and size guards; the "
        "ledger hash chain is verified (verified flag). Provide url, agent_id, "
        "and optional pack_id."
    ),
)
def import_pack_by_ref(req: ImportByRefRequest) -> ImportByRefResponse:
    """Fetch, verify, and import a pack referenced by a remote HTTPS URL."""
    try:
        raw = pack_import.fetch_pack_bytes(req.url, _max_import_bytes())
    except ImportRefError as exc:
        raise HTTPException(status_code=exc.status_code, detail=exc.detail)
    return _import_pack_from_bytes(
        raw, req.agent_id, req.pack_id, "URL did not return valid JSON."
    )


@router.post(
    "/import-from-upload-ref",
    response_model=ImportByRefResponse,
    operation_id="importMemoryPackFromUploadRef",
    include_in_schema=False,
)
def import_pack_from_upload_ref(
    req: ImportFromUploadRefRequest,
    auth: UploadAuthContext = Depends(require_upload_or_token),
) -> ImportByRefResponse:
    """Import a pre-built .smp.json pack from a staged upload.

    Binary packs cannot travel through the LLM/GPT Actions, so the browser first
    stages the file via /api/uploads, then calls this with the upload_id. Two
    callers, resolved by auth:
    * The /import page (master key) imports into agent_id's persistent vault.
    * The keyless /u/{token} import page (X-Upload-Token) imports into the
      link's private, unguessable, TTL-expired namespace, so the user can query
      only their own packs for the life of the link (tenant isolation, no
      persistence). Hidden from OpenAPI; staged bytes are single-use.
    """
    claim = auth.claim if auth.mode == "token" else None
    # Opportunistic hygiene on every import: remove expired ephemeral packs from
    # earlier links so the server does not accumulate stale one-time uploads.
    jobs_store.sweep_expired_quietly()
    if claim is not None and (getattr(claim, "mode", "build") or "build") != "import":
        # A build-mode one-time token must not be repurposed to import packs.
        raise HTTPException(
            status_code=403, detail="This link does not allow importing packs."
        )

    # In token mode the pack is re-homed under the claim's ephemeral namespace,
    # never the caller-supplied agent_id, so links cannot reach each other.
    target_agent_id = (
        (claim.import_agent_id if claim is not None else None) or req.agent_id
    )

    store = storage.get_storage()
    record = storage.load_upload_record(req.upload_id)
    if record is None or not store.exists(req.upload_id):
        raise HTTPException(status_code=404, detail="Staged upload not found.")

    raw = store.open(req.upload_id)
    if raw is None:
        raise HTTPException(status_code=404, detail="Staged upload not found.")
    try:
        result = _import_pack_from_bytes(
            raw, target_agent_id, req.pack_id, "Uploaded file is not valid JSON."
        )
    finally:
        # Imported (or rejected) staged bytes are single-use; clean up.
        try:
            store.delete(req.upload_id)
        except Exception:  # noqa: BLE001 - cleanup is best-effort.
            pass

    # Ephemeral (token) imports get a temporary-retention job so the existing
    # cleanup sweep removes the pack file once the link expires. Never vaulted.
    if claim is not None:
        _register_ephemeral_import(claim, target_agent_id, result)

    return result


def _register_ephemeral_import(
    claim, agent_id: str, result: ImportByRefResponse
) -> None:
    """Bind an imported pack to a temp job and record it on the import claim.

    The temp JobRecord points cleanup_expired_temp_jobs at the imported pack
    file, so it is deleted after the link's TTL (never persisted to the vault).
    Only safe summary fields (no secret content) are stored on the claim.
    """
    found = pack_io.find_pack_by_id(agent_id, result.pack_id)
    pack_rel = pack_io.relpath_from_root(found) if found is not None else None
    job = JobRecord(
        job_id=uuid.uuid4().hex,
        agent_id=agent_id,
        pack_id=result.pack_id,
        status=JobStatus.COMPLETED,
        retention_mode=RetentionMode.SESSION,
        expires_at=(
            claim.expires_at
            or jobs_store.compute_expires_at(RetentionMode.SESSION)
        ),
        pack_persisted=True,
        entry_count=result.entry_count,
        classification_counts=dict(result.classification_summary),
        pack_path=pack_rel,
    )
    jobs_store.save_job(job)
    upload_links.record_import(
        claim,
        {
            "agent_id": agent_id,
            "pack_id": result.pack_id,
            "entry_count": result.entry_count,
            "classifications": dict(result.classification_summary),
            "verified": bool(result.verified),
            "job_id": job.job_id,
        },
    )


def _looks_japanese(text: str) -> bool:
    """Heuristic: does the text contain CJK/Japanese characters?"""
    return bool(_CJK_RE.search(text or ""))


def _rows_from_xlsx(data: bytes) -> List[str]:
    """Extract non-empty, header-labeled records from an uploaded .xlsx."""
    import openpyxl  # local import keeps the module import light

    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    records: List[str] = []
    try:
        for ws in wb.worksheets:
            headers: Optional[List[str]] = None
            for r_index, row in enumerate(ws.iter_rows(values_only=True)):
                non_empty = [c for c in row if c is not None and str(c).strip()]
                if not non_empty:
                    continue
                if r_index == 0 and headers is None:
                    headers = [
                        str(c).strip() if c is not None else "" for c in row
                    ]
                    continue
                parts: List[str] = []
                for idx, cell in enumerate(row):
                    if cell is None:
                        continue
                    value = str(cell).strip()
                    if not value:
                        continue
                    if headers and idx < len(headers) and headers[idx]:
                        parts.append(f"{headers[idx]}: {value}")
                    else:
                        parts.append(value)
                record = "\n".join(parts)
                if record:
                    records.append(record)
    finally:
        wb.close()
    return records


# Image extensions handled via OCR (Tesseract).
_OCR_IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp")


# Generic label headers for a 2-column key/value table. When the first row is
# made entirely of these, it is a header row and skipped (the real key/value
# pairs live in the data rows, e.g. "Budget cap: USD 420,000").
_TABLE_LABEL_HEADERS = frozenset(
    {
        "field",
        "fields",
        "value",
        "values",
        "item",
        "items",
        "key",
        "keys",
        "attribute",
        "attributes",
        "property",
        "properties",
        "name",
        "description",
        "detail",
        "details",
        "項目",
        "内容",
        "値",
        "名称",
        "属性",
        "説明",
    }
)


def _rows_from_xls(data: bytes) -> List[str]:
    """Extract header-labeled records from a legacy binary .xls via xlrd.

    Mirrors :func:`_rows_from_xlsx`: the first non-empty row of each sheet is
    treated as headers and subsequent rows become ``header: value`` records.
    """
    import xlrd  # local import keeps the module import light

    book = xlrd.open_workbook(file_contents=data)
    records: List[str] = []
    for sheet in book.sheets():
        headers: Optional[List[str]] = None
        header_row_seen = False
        for r_index in range(sheet.nrows):
            row = sheet.row_values(r_index)
            cells = ["" if c is None else str(c).strip() for c in row]
            if not any(cells):
                continue
            if not header_row_seen:
                headers = cells
                header_row_seen = True
                continue
            parts: List[str] = []
            for idx, value in enumerate(cells):
                if not value:
                    continue
                if headers and idx < len(headers) and headers[idx]:
                    parts.append(f"{headers[idx]}: {value}")
                else:
                    parts.append(value)
            record = "\n".join(parts)
            if record:
                records.append(record)
    return records


def _records_from_table(rows: List[List[str]]) -> List[str]:
    """Turn a table (docx/pptx) into retrieval-friendly ``header: value`` records.

    Mirrors :func:`_rows_from_xlsx` so table content embeds as well as spreadsheet
    content:

    * A **2-column** table is treated as a key/value sheet: each data row becomes
      ``"<col0>: <col1>"`` (e.g. ``"Budget cap: USD 420,000"``). An obvious generic
      label header row (``Field``/``Value``/``項目``/``内容`` ...) is skipped.
    * A **wider** table uses its first non-empty row as column headers and emits
      ``"<header>: <value>"`` per non-empty cell, newline-joined per row (columns
      whose header cell is empty contribute the bare value).
    """
    cleaned = [
        [(c or "").strip() for c in row]
        for row in rows
        if any((c or "").strip() for c in row)
    ]
    if not cleaned:
        return []

    ncols = max(len(r) for r in cleaned)
    records: List[str] = []

    if ncols <= 2:
        start = 0
        first = cleaned[0]
        labels = [c for c in first[:2] if c]
        if labels and all(c.lower() in _TABLE_LABEL_HEADERS for c in labels):
            start = 1
        for row in cleaned[start:]:
            key = row[0] if len(row) > 0 else ""
            val = row[1] if len(row) > 1 else ""
            if key and val:
                records.append(f"{key}: {val}")
            elif key or val:
                records.append(key or val)
        return records

    headers = cleaned[0]
    for row in cleaned[1:]:
        parts: List[str] = []
        for idx, value in enumerate(row):
            if not value:
                continue
            if idx < len(headers) and headers[idx]:
                parts.append(f"{headers[idx]}: {value}")
            else:
                parts.append(value)
        record = "\n".join(parts)
        if record:
            records.append(record)
    return records


def _records_from_docx(data: bytes) -> List[str]:
    """Extract paragraph and table text from a Word .docx file.

    Paragraphs become one record each; tables are rendered as retrieval-friendly
    ``header: value`` records via :func:`_records_from_table`.
    """
    import docx  # local import (python-docx)

    document = docx.Document(io.BytesIO(data))
    records: List[str] = []
    for para in document.paragraphs:
        text = (para.text or "").strip()
        if text:
            records.append(text)
    for table in document.tables:
        rows = [[(cell.text or "") for cell in row.cells] for row in table.rows]
        records.extend(_records_from_table(rows))
    return records


def _records_from_pptx(data: bytes) -> List[str]:
    """Extract shape text, tables, and speaker notes from a PowerPoint .pptx file.

    Each slide contributes one record combining its non-table shape text and
    notes (so provenance stays slide-scoped); table shapes are emitted separately
    as retrieval-friendly ``header: value`` records via :func:`_records_from_table`.
    """
    from pptx import Presentation  # local import (python-pptx)

    presentation = Presentation(io.BytesIO(data))
    records: List[str] = []
    for slide in presentation.slides:
        lines: List[str] = []
        table_records: List[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False) and shape.has_table:
                rows = [
                    [(cell.text or "") for cell in row.cells]
                    for row in shape.table.rows
                ]
                table_records.extend(_records_from_table(rows))
                continue
            if shape.has_text_frame:
                text = (shape.text_frame.text or "").strip()
                if text:
                    lines.append(text)
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                lines.append(f"Notes: {notes}")
        record = "\n".join(lines).strip()
        if record:
            records.append(record)
        records.extend(table_records)
    return records


def _pdf_text_pages(data: bytes) -> List[str]:
    """Return the extracted text of each PDF page (empty string when none)."""
    from pypdf import PdfReader  # local import

    reader = PdfReader(io.BytesIO(data))
    pages: List[str] = []
    for page in reader.pages:
        try:
            pages.append((page.extract_text() or "").strip())
        except Exception:  # pragma: no cover - defensive (malformed page)
            pages.append("")
    return pages


def _ocr_images(images) -> List[str]:
    """OCR a list of PIL images with Tesseract (Japanese + English by default).

    Isolated so tests can monkeypatch it without a real Tesseract binary. Raises
    :class:`UploadProcessingError` when the Tesseract engine is unavailable.
    """
    try:
        import pytesseract  # local import
    except ImportError as exc:  # pragma: no cover - dependency guard
        raise UploadProcessingError(
            500, "OCR support (pytesseract) is not installed on the server."
        ) from exc

    lang = settings.safe_memory_ocr_languages or "eng"
    texts: List[str] = []
    for image in images:
        try:
            text = pytesseract.image_to_string(image, lang=lang)
        except pytesseract.TesseractNotFoundError as exc:
            raise UploadProcessingError(
                500,
                "OCR engine (Tesseract) is not available on the server.",
            ) from exc
        except Exception as exc:
            raise UploadProcessingError(400, f"OCR failed: {exc}") from exc
        cleaned = (text or "").strip()
        if cleaned:
            texts.append(cleaned)
    return texts


def _ocr_pdf(data: bytes) -> List[str]:
    """Rasterize a PDF (capped page count) and OCR each page image."""
    try:
        from pdf2image import convert_from_bytes  # local import
    except ImportError as exc:  # pragma: no cover - dependency guard
        raise UploadProcessingError(
            500, "OCR support (pdf2image) is not installed on the server."
        ) from exc

    max_pages = max(1, int(settings.safe_memory_max_ocr_pages))
    dpi = max(72, int(settings.safe_memory_ocr_dpi))
    try:
        images = convert_from_bytes(
            data, dpi=dpi, first_page=1, last_page=max_pages
        )
    except UploadProcessingError:
        raise
    except Exception as exc:
        raise UploadProcessingError(
            400, f"Could not rasterize the PDF for OCR: {exc}"
        ) from exc
    return _ocr_images(images)


def _records_from_pdf(data: bytes) -> List[str]:
    """Extract records from a PDF, falling back to OCR for scanned pages.

    First tries fast text extraction with pypdf. When the embedded text is empty
    or implausibly sparse for the page count (a scanned/image PDF), rasterize and
    OCR instead.
    """
    try:
        pages = _pdf_text_pages(data)
    except UploadProcessingError:
        raise
    except Exception as exc:
        raise UploadProcessingError(400, f"Could not read the PDF: {exc}") from exc

    combined = "\n".join(p for p in pages if p).strip()
    page_count = max(1, len(pages))
    # Heuristic: treat as scanned when there is essentially no embedded text
    # (fewer than ~10 characters per page on average).
    if len(combined) < 10 * page_count:
        ocr_texts = _ocr_pdf(data)
        if ocr_texts:
            return ocr_texts
    records = [p for p in pages if p]
    return records


def _records_from_image(data: bytes) -> List[str]:
    """OCR a single uploaded image into text records."""
    try:
        from PIL import Image  # local import
    except ImportError as exc:  # pragma: no cover - dependency guard
        raise UploadProcessingError(
            500, "Image support (Pillow) is not installed on the server."
        ) from exc
    try:
        image = Image.open(io.BytesIO(data))
        image.load()
    except UploadProcessingError:
        raise
    except Exception as exc:
        raise UploadProcessingError(400, f"Could not read the image: {exc}") from exc
    return _ocr_images([image])


def _build_upload_specs(
    raw_texts: List[str],
    source_language: Optional[str],
    canonical_language: str,
    source: str,
    sources: Optional[List[str]] = None,
    *,
    translate: bool = False,
    apply_default_classification: bool = False,
    default_classification: Classification = Classification.INTERNAL,
) -> List[Dict[str, Any]]:
    """Turn raw text records into bilingual entry specs, translating if asked.

    Translation is OFF by default (``translate=False``): canonical_text is the
    original text verbatim and no Qwen chat calls are made. When
    ``translate=True``, only rows that actually contain Japanese characters are
    sent to Qwen in a single batched call (falling back to per-item translation
    only if a batch misaligns).

    When ``apply_default_classification`` is True (the caller explicitly supplied
    a classification), every entry is stamped with ``default_classification`` and
    the expensive per-entry LLM classification is skipped -- except that obvious
    secrets are still promoted to SECRET via a keyword heuristic (no LLM call).

    ``sources`` (optional) is a per-record list of provenance source names
    (used for folder-ZIP imports so each entry records its origin filename);
    when omitted, every entry uses the single ``source`` value.
    """
    cleaned = [(raw or "").strip() for raw in raw_texts]

    # Decide per row whether translation is needed, preserving original order.
    # Translation only happens when explicitly requested AND the row actually
    # contains Japanese characters; structured ASCII data (IDs, numbers, FX
    # tickers) is never sent to the model.
    if translate:
        needs: List[bool] = [bool(raw) and _looks_japanese(raw) for raw in cleaned]
    else:
        needs = [False] * len(cleaned)

    # Batch-translate only the rows that need it.
    to_translate = [raw for raw, flag in zip(cleaned, needs) if flag]
    translated_iter = iter(
        normalize_accounting_batch(to_translate, batch_size=settings.translation_batch_size)
        if to_translate
        else []
    )

    def _fast_classification(text: str) -> Classification:
        """Apply the supplied default, promoting only obvious secrets."""
        if qwen_client.has_secret_terms(text):
            return Classification.SECRET
        return default_classification

    specs: List[Dict[str, Any]] = []
    for index, (raw, needs_translation) in enumerate(zip(cleaned, needs)):
        if not raw:
            continue
        entry_source = sources[index] if sources else source
        if needs_translation:
            canonical = next(translated_iter)
            note = (
                "Auto-translated to English via Qwen."
                if not canonical.startswith("[UNTRANSLATED FALLBACK]")
                else "Qwen unavailable; canonical text is a labeled fallback."
            )
            spec: Dict[str, Any] = {
                "text": canonical,
                "original_text": raw,
                "canonical_text": canonical,
                "source_language": source_language or "ja",
                "canonical_language": canonical_language,
                "translation_note": note,
                "source": entry_source,
            }
        else:
            spec = {
                "text": raw,
                "canonical_text": raw,
                "source_language": source_language or canonical_language,
                "canonical_language": canonical_language,
                "source": entry_source,
            }
        # Stamp an explicit classification so build_pack_from_entries skips the
        # per-entry LLM classification (major speedup for large imports).
        if apply_default_classification:
            classify_text = f"{raw}\n{spec.get('canonical_text', '')}"
            spec["classification"] = _fast_classification(classify_text)
        specs.append(spec)
    return specs


class UploadProcessingError(Exception):
    """Raised inside the shared upload pipeline with an HTTP-friendly status."""

    def __init__(self, status_code: int, detail: str) -> None:
        self.status_code = status_code
        self.detail = detail
        super().__init__(detail)


def _parse_upload_records(filename: str, data: bytes) -> List[str]:
    """Extract text records from an uploaded file by extension."""
    name = (filename or "").lower()
    if name.endswith(".xlsx"):
        try:
            return _rows_from_xlsx(data)
        except Exception:
            raise UploadProcessingError(400, "Could not read the .xlsx file.")
    if name.endswith(".xls"):
        try:
            return _rows_from_xls(data)
        except UploadProcessingError:
            raise
        except Exception:
            raise UploadProcessingError(400, "Could not read the .xls file.")
    if name.endswith(".docx"):
        try:
            return _records_from_docx(data)
        except UploadProcessingError:
            raise
        except Exception:
            raise UploadProcessingError(400, "Could not read the .docx file.")
    if name.endswith(".pptx"):
        try:
            return _records_from_pptx(data)
        except UploadProcessingError:
            raise
        except Exception:
            raise UploadProcessingError(400, "Could not read the .pptx file.")
    if name.endswith(".pdf"):
        return _records_from_pdf(data)
    if name.endswith(_OCR_IMAGE_EXTS):
        return _records_from_image(data)
    if name.endswith(".csv") or name.endswith(".tsv"):
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            text = data.decode("utf-8", errors="replace")
        rows = [line.strip() for line in text.splitlines() if line.strip()]
        return rows or chunk_text(text)
    if name.endswith(".txt") or name.endswith(".md") or not name:
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            text = data.decode("utf-8", errors="replace")
        return chunk_text(text)
    if name.endswith(".json"):
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            text = data.decode("utf-8", errors="replace")
        return _records_from_json(text)
    raise UploadProcessingError(
        415,
        "Unsupported file type. Use .txt, .md, .csv, .tsv, .json, .xlsx, .xls, "
        ".docx, .pptx, .pdf, or an image (.png/.jpg/.jpeg/.tiff/.bmp/.webp).",
    )


def _records_from_json(text: str) -> List[str]:
    """Ingest a .json file as text records.

    A JSON array becomes one record per item; an object becomes ``key: value``
    lines; anything else falls back to chunking the raw text.
    """
    try:
        parsed = json.loads(text)
    except Exception:
        return chunk_text(text)

    def _stringify(value) -> str:
        if isinstance(value, str):
            return value.strip()
        return json.dumps(value, ensure_ascii=False)

    records: List[str] = []
    if isinstance(parsed, list):
        for item in parsed:
            rec = _stringify(item)
            if rec.strip():
                records.append(rec)
    elif isinstance(parsed, dict):
        for key, value in parsed.items():
            records.append(f"{key}: {_stringify(value)}")
    else:
        return chunk_text(text)
    return records or chunk_text(text)


# Extensions we can actually ingest when pulling files out of a folder ZIP.
_SUPPORTED_MEMBER_EXTS = (
    ".xlsx",
    ".xls",
    ".csv",
    ".tsv",
    ".txt",
    ".md",
    ".json",
    ".docx",
    ".pptx",
    ".pdf",
) + _OCR_IMAGE_EXTS

# Known extensions we recognize but cannot ingest yet -> recorded as unsupported.
# (docx/pptx/pdf/xls and images are now handled, so nothing common remains here.)
_UNSUPPORTED_KNOWN_EXTS = (".doc", ".ppt", ".key", ".pages", ".numbers")


def _member_is_unsafe(name: str) -> bool:
    """Zip-slip guard: reject absolute paths, drive letters, or ``..`` segments."""
    normalized = (name or "").replace("\\", "/")
    if normalized.startswith("/"):
        return True
    if len(normalized) >= 2 and normalized[1] == ":":  # Windows drive (e.g. C:)
        return True
    return ".." in normalized.split("/")


def _is_supported_member(name: str) -> bool:
    base = (name or "").replace("\\", "/").rsplit("/", 1)[-1]
    return base.lower().endswith(_SUPPORTED_MEMBER_EXTS)


def _is_folder_zip(data: bytes) -> bool:
    """True when ``data`` is a plain ZIP archive of files (a shared folder).

    Office files (.xlsx/.docx/.pptx) are themselves ZIPs, so they are excluded
    by looking for the OPC marker ``[Content_Types].xml`` at the archive root --
    when present the payload is a single Office document, not a folder.
    """
    if not zipfile.is_zipfile(io.BytesIO(data)):
        return False
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            names = archive.namelist()
    except zipfile.BadZipFile:
        return False
    if any(n == "[Content_Types].xml" for n in names):
        return False
    return True


def _extract_folder_zip(data: bytes):
    """Extract supported files from a folder ZIP into merged records.

    Returns ``(records, sources, warnings, unsupported)`` where ``sources[i]`` is
    the origin filename of ``records[i]`` (kept as provenance) and ``unsupported``
    is a list of ``{"filename", "reason"}`` for recognized-but-unsupported members
    (e.g. legacy .doc/.ppt). Enforces zip-slip safety and folder file-count /
    total-uncompressed-size caps (zip-bomb guard). Hidden / __MACOSX entries are
    skipped; unreadable supported files are skipped with a warning.
    """
    records: List[str] = []
    sources: List[str] = []
    warnings: List[str] = []
    unsupported: List[Dict[str, str]] = []
    max_bytes = _max_upload_bytes()
    max_files = settings.safe_memory_max_folder_files
    max_total = settings.safe_memory_max_folder_total_size_mb * 1024 * 1024

    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        infos = archive.infolist()
        total_uncompressed = sum(info.file_size for info in infos)
        if total_uncompressed > min(max_bytes, max_total):
            raise UploadProcessingError(
                413,
                "Archive contents exceed the "
                f"{settings.safe_memory_max_folder_total_size_mb} MB folder limit.",
            )
        file_count = 0
        for info in infos:
            name = info.filename
            if info.is_dir():
                continue
            # Zip-slip: validate every entry before deciding to skip it.
            if _member_is_unsafe(name):
                raise UploadProcessingError(400, "Unsafe path in archive entry.")
            base = name.replace("\\", "/").rsplit("/", 1)[-1]
            if name.startswith("__MACOSX") or base.startswith("."):
                continue
            lower = base.lower()
            if lower.endswith(_UNSUPPORTED_KNOWN_EXTS):
                unsupported.append(
                    {"filename": base, "reason": "unsupported file type"}
                )
                continue
            if not _is_supported_member(name):
                continue
            file_count += 1
            if file_count > max_files:
                raise UploadProcessingError(
                    413,
                    f"Folder archive exceeds the {max_files} file limit.",
                )
            try:
                member_bytes = archive.read(info)
                member_records = _parse_upload_records(base, member_bytes)
            except UploadProcessingError as exc:
                warnings.append(f"Skipped {base}: {exc.detail}")
                unsupported.append({"filename": base, "reason": exc.detail})
                continue
            except Exception:  # pragma: no cover - defensive (corrupt member)
                warnings.append(f"Skipped {base}: could not read file.")
                continue
            for record in member_records:
                records.append(record)
                sources.append(base)
    return records, sources, warnings, unsupported


def _collect_upload_records(filename: str, data: bytes):
    """Return ``(records, sources, warnings, unsupported, input_type)``.

    For a folder ZIP, records from every supported member are merged into one
    list and ``sources`` carries each record's origin filename. For a single
    file, ``sources`` is ``None`` (all entries share the default source).
    """
    if _is_folder_zip(data):
        records, sources, warnings, unsupported = _extract_folder_zip(data)
        if not records:
            reason = (
                "No supported files (.xlsx/.xls/.csv/.tsv/.txt/.md/.json/.docx/"
                ".pptx/.pdf/images) found in the folder archive."
            )
            raise UploadProcessingError(400, reason)
        return records, sources, warnings, unsupported, "folder"
    return _parse_upload_records(filename, data), None, [], [], "file"


def _run_pack_build(
    *,
    job_id: str,
    agent_id: str,
    pack_id: str,
    title: str,
    data: bytes,
    filename: str,
    source_language: Optional[str],
    canonical_language: str,
    default_classification: Classification,
    retention_mode: RetentionMode,
    debug_keep_upload: bool,
    return_download_url: bool = True,
    delete_source_after_processing: bool = True,
    translate: bool = False,
    apply_default_classification: bool = False,
) -> "tuple[JobRecord, str]":
    """Shared pipeline for BOTH the multipart and staged-ref upload endpoints.

    Parses records (single file or folder ZIP), batch-translates, builds the pack
    (temp vs vault by retention_mode), optionally hands the pack off to Alibaba
    OSS (private bucket + signed URL) when configured, deletes working files
    unless kept, and returns the completed (unsaved) :class:`JobRecord` plus the
    audit path. Raises :class:`UploadProcessingError` on any recoverable failure.
    """
    if not data:
        raise UploadProcessingError(400, "Uploaded file is empty.")
    if len(data) > _max_upload_bytes():
        raise UploadProcessingError(
            413,
            f"Upload exceeds the {settings.safe_memory_max_upload_mb} MB limit.",
        )

    records, per_record_sources, extract_warnings, unsupported, input_type = (
        _collect_upload_records(filename, data)
    )
    if not records:
        raise UploadProcessingError(400, "No usable content in the upload.")

    warnings: List[str] = list(extract_warnings)
    # The raw upload is needed only during processing. Persist it to a per-job
    # working directory so cleanup is explicit and testable.
    working_rel = jobs_store.write_working_upload(job_id, filename or "upload", data)

    specs = _build_upload_specs(
        records,
        source_language,
        canonical_language,
        source="upload",
        sources=per_record_sources,
        translate=translate,
        apply_default_classification=apply_default_classification,
        default_classification=default_classification,
    )

    # Choose where the generated pack lives based on retention mode.
    if retention_mode == RetentionMode.SERVER_VAULT:
        target = pack_io.pack_target_path(agent_id, pack_id, default_classification)
        pack_persisted = True
    else:
        target = pack_io.temp_pack_target_path(agent_id, pack_id, job_id)
        pack_persisted = False

    try:
        pack, saved_path, audit_path = build_pack_from_entries(
            agent_id=agent_id,
            pack_id=pack_id,
            title=title,
            entries=specs,
            default_classification=default_classification,
            method="upload",
            target_path=target,
        )
    except ValueError as exc:
        jobs_store.delete_working_files(
            JobRecord(
                job_id=job_id,
                agent_id=agent_id,
                pack_id=pack_id,
                working_dir=working_rel,
            )
        )
        raise UploadProcessingError(400, str(exc))

    counts = dict(Counter(e.classification.value for e in pack.entries))
    pack_rel = pack_io.relpath_from_root(saved_path)

    # server_vault packs are catalog-visible; temporary modes are not.
    catalog_visible = retention_mode == RetentionMode.SERVER_VAULT

    if retention_mode in (RetentionMode.SESSION, RetentionMode.PROCESS_AND_RETURN):
        download_url = f"/api/jobs/{job_id}/download"
    else:
        download_url = None
    expires_at = jobs_store.compute_expires_at(retention_mode)

    # ---- Optional Alibaba OSS handoff (private bucket + signed URL) ----
    oss_object_key: Optional[str] = None
    oss_uploaded = False
    want_oss = return_download_url and (
        retention_mode
        in (RetentionMode.SESSION, RetentionMode.PROCESS_AND_RETURN)
        or (retention_mode == RetentionMode.SERVER_VAULT and return_download_url)
    )
    if want_oss and oss_storage.is_enabled():
        try:
            object_key = f"{settings.oss_export_prefix}{job_id}/{pack_id}.smp.json"
            oss_storage.upload_file(
                saved_path, object_key, content_type="application/json"
            )
            oss_object_key = object_key
            oss_uploaded = True
            # download_url becomes a stable token URL below (never the raw
            # signed URL, whose base64 signature can be corrupted in transit).
            download_url = None
        except Exception as exc:  # pragma: no cover - defensive; keep local fallback
            warnings.append(
                f"OSS upload failed ({type(exc).__name__}); using local download."
            )
            if retention_mode in (
                RetentionMode.SESSION,
                RetentionMode.PROCESS_AND_RETURN,
            ):
                download_url = f"/api/jobs/{job_id}/download"

    # ---- Stable, signature-free download token ----
    # Hand GPT/browsers a URL that never carries a base64 OSS signature. The
    # token URL streams the local pack when present, else 307-redirects to a
    # freshly signed OSS URL. Minted whenever a download is expected (temporary
    # modes), a pack was uploaded to OSS, or the pack is persisted locally
    # (server_vault) -- so a COMPLETED job ALWAYS carries a download_url,
    # independent of retention_mode or whether OSS is enabled/succeeded.
    download_token: Optional[str] = None
    if download_url is not None or oss_object_key is not None or pack_persisted:
        download_token = export_links.create_export_link(
            agent_id,
            pack_rel,
            oss_object_key=oss_object_key,
            job_id=job_id,
        )
        download_url = _absolute_url(f"/api/packs/dl/{download_token}")

    keep_raw = debug_keep_upload or not delete_source_after_processing
    raw_deleted = False
    if keep_raw:
        warnings.append("raw upload retained in working dir.")
    else:
        jobs_store.delete_working_files(
            JobRecord(
                job_id=job_id,
                agent_id=agent_id,
                pack_id=pack_id,
                working_dir=working_rel,
            )
        )
        raw_deleted = True

    job = JobRecord(
        job_id=job_id,
        agent_id=agent_id,
        pack_id=pack.manifest.pack_id,
        status=JobStatus.COMPLETED,
        retention_mode=retention_mode,
        expires_at=expires_at,
        raw_upload_deleted=raw_deleted,
        working_files_deleted=raw_deleted,
        pack_persisted=pack_persisted,
        download_url=download_url,
        download_token=download_token,
        entry_count=len(pack.entries),
        classification_counts=counts,
        warnings=warnings,
        input_type=input_type,
        catalog_visible=catalog_visible,
        oss_export_uploaded=oss_uploaded,
        oss_object_key=oss_object_key,
        unsupported_files=unsupported,
        pack_path=pack_rel,
        working_dir=None if raw_deleted else working_rel,
    )
    return job, pack_io.relpath_from_root(audit_path)


def _coerce_retention_mode(value) -> RetentionMode:
    """Coerce a string/enum retention_mode, raising HTTP 400 when invalid."""
    if isinstance(value, RetentionMode):
        return value
    try:
        return RetentionMode(value or "process_and_return")
    except ValueError:
        valid = ", ".join(m.value for m in RetentionMode)
        raise HTTPException(
            status_code=400, detail=f"Invalid retention_mode. Use one of: {valid}."
        )


def _coerce_classification(value) -> Classification:
    """Coerce a string/enum classification, defaulting to INTERNAL when invalid."""
    if isinstance(value, Classification):
        return value
    try:
        return Classification(value)
    except (ValueError, TypeError):
        return Classification.INTERNAL


def _upload_build_response(
    job: JobRecord, debug: bool, audit_rel: Optional[str]
) -> UploadBuildResponse:
    resp = UploadBuildResponse(
        pack_id=job.pack_id,
        entry_count=job.entry_count,
        classification_counts=job.classification_counts,
        job_id=job.job_id,
        status=job.status.value,
        retention_mode=job.retention_mode.value,
        expires_at=job.expires_at,
        download_url=job.download_url,
        warnings=job.warnings,
    )
    if debug:
        resp.pack_path = job.pack_path
        resp.audit_path = audit_rel
    return resp


@router.post(
    "/build-from-upload",
    response_model=UploadBuildResponse,
    operation_id="buildMemoryPackFromUpload",
    include_in_schema=False,
    summary="Build a Safe Memory Pack from an uploaded file",
    description=(
        "Upload a document (.txt/.md/.csv/.tsv/.json/.xlsx/.xls/.docx/.pptx/.pdf) "
        "or image and build a Safe Memory Pack. Scanned PDFs and images are OCR'd "
        "(Japanese+English). Raw uploads are session-scoped and deleted after "
        "processing. retention_mode sets temporary vs vault storage. Returns job "
        "metadata."
    ),
)
async def build_pack_from_upload(
    agent_id: str = Form(...),
    pack_id: str = Form(...),
    title: str = Form(...),
    file: UploadFile = File(...),
    source_language: Optional[str] = Form(default=None),
    canonical_language: str = Form(default="en"),
    default_classification: Optional[Classification] = Form(default=None),
    retention_mode: RetentionMode = Form(default=RetentionMode.PROCESS_AND_RETURN),
    debug_keep_upload: bool = Form(default=False),
    debug: bool = Form(default=False),
    translate: bool = Form(default=False),
) -> UploadBuildResponse:
    """Build a Safe Memory Pack from an uploaded text/markdown/Excel file."""
    data = await file.read()
    # When the caller explicitly names a classification, apply it to every entry
    # and skip the per-entry LLM classification (big speedup); otherwise keep the
    # backward-compatible per-entry LLM behaviour with an INTERNAL default.
    apply_default_classification = default_classification is not None
    effective_classification = default_classification or Classification.INTERNAL
    job_id = uuid.uuid4().hex
    try:
        job, audit_rel = _run_pack_build(
            job_id=job_id,
            agent_id=agent_id,
            pack_id=pack_id,
            title=title,
            data=data,
            filename=file.filename or "upload",
            source_language=source_language,
            canonical_language=canonical_language,
            default_classification=effective_classification,
            retention_mode=retention_mode,
            debug_keep_upload=debug_keep_upload,
            translate=translate,
            apply_default_classification=apply_default_classification,
        )
    except UploadProcessingError as exc:
        raise HTTPException(status_code=exc.status_code, detail=exc.detail)

    jobs_store.save_job(job)
    return _upload_build_response(job, debug, audit_rel)


def _mark_ref_job_failed(
    job_id: str, req: BuildFromUploadRefRequest, message: str
) -> None:
    """Update (or create) the job record as FAILED with an error warning."""
    job = jobs_store.load_job(job_id)
    if job is None:
        job = JobRecord(
            job_id=job_id,
            agent_id=req.agent_id,
            pack_id=req.pack_id,
            retention_mode=_safe_retention(req.retention_mode),
        )
    job.status = JobStatus.FAILED
    job.warnings = list(job.warnings) + [f"Processing failed: {message}"]
    jobs_store.save_job(job)


def _safe_retention(value) -> RetentionMode:
    try:
        return value if isinstance(value, RetentionMode) else RetentionMode(value)
    except ValueError:
        return RetentionMode.PROCESS_AND_RETURN


def _process_ref_job(
    job_id: str,
    upload_id: str,
    req: BuildFromUploadRefRequest,
    apply_default_classification: bool = False,
) -> None:
    """Background worker: build the pack from staged bytes and finalize the job."""
    store = storage.get_storage()
    try:
        record = storage.load_upload_record(upload_id)
        data = store.open(upload_id)
        if record is None or data is None:
            raise UploadProcessingError(404, "Staged upload content not found.")

        job, _audit = _run_pack_build(
            job_id=job_id,
            agent_id=req.agent_id,
            pack_id=req.pack_id,
            title=req.title,
            data=data,
            filename=record.filename,
            source_language=req.source_language,
            canonical_language=req.canonical_language,
            default_classification=req.default_classification,
            retention_mode=_safe_retention(req.retention_mode),
            debug_keep_upload=req.debug_keep_upload,
            return_download_url=getattr(req, "return_download_url", True),
            delete_source_after_processing=getattr(
                req, "delete_source_after_processing", True
            ),
            translate=getattr(req, "translate", False),
            apply_default_classification=apply_default_classification,
        )
        jobs_store.save_job(job)

        # Remove the staged upload after processing unless kept for debugging.
        if req.debug_keep_upload:
            record.status = UploadStatus.CONSUMED
            storage.save_upload_record(record)
        else:
            store.delete(upload_id)
    except UploadProcessingError as exc:
        _mark_ref_job_failed(job_id, req, exc.detail)
    except Exception as exc:  # pragma: no cover - defensive
        _mark_ref_job_failed(job_id, req, f"Unexpected error ({type(exc).__name__}).")


# Strong references to in-flight background build tasks so they are not
# garbage-collected when a bounded-wait endpoint returns early (PROCESSING).
_INFLIGHT_BUILD_TASKS: "set[asyncio.Task]" = set()


async def _run_ref_job_bounded(
    job_id: str,
    upload_id: str,
    ref_req: "BuildFromUploadRefRequest",
    apply_default_classification: bool,
    wait_seconds: Optional[float] = None,
) -> Optional[JobRecord]:
    """Run the ingest job as an awaitable task, waiting up to a bound for it.

    The blocking :func:`_process_ref_job` runs in a worker thread. We wait up to
    ``wait_seconds`` for it to reach a terminal state. If it finishes in time the
    finalized :class:`JobRecord` is returned; if the wait elapses the task is
    **not** cancelled (``asyncio.shield`` keeps it running in the background) and
    the current (still PROCESSING) job is returned so the caller can poll.
    """
    if wait_seconds is None:
        wait_seconds = float(settings.safe_memory_sync_build_wait_seconds)

    task = asyncio.create_task(
        asyncio.to_thread(
            _process_ref_job,
            job_id,
            upload_id,
            ref_req,
            apply_default_classification,
        )
    )
    _INFLIGHT_BUILD_TASKS.add(task)
    task.add_done_callback(_INFLIGHT_BUILD_TASKS.discard)

    try:
        await asyncio.wait_for(asyncio.shield(task), timeout=max(0.0, wait_seconds))
    except asyncio.TimeoutError:
        # Leave the shielded task running; the caller falls back to polling.
        return jobs_store.load_job(job_id)
    return jobs_store.load_job(job_id)


async def _build_pack_from_url_impl(
    req: "BuildFromUrlRequest", wait_seconds: Optional[float] = None
):
    """Fetch a file from a URL and build a pack with bounded-synchronous waiting.

    Always returns a :class:`JobResponse`. When the build finishes within the
    wait window (including early folder/HTML failures) it is a terminal job with
    ``download_url``; otherwise it is a ``PROCESSING`` job carrying ``upload_id``
    so the caller polls GET /api/jobs/{job_id}.
    """
    retention = _coerce_retention_mode(req.retention_mode)

    # Some folder links (e.g. Google Drive folders, or SharePoint/OneDrive
    # folders identified by the `:f:` marker) have no anonymous ZIP download;
    # fail the job clearly instead of fetching an HTML page.
    folder_err = url_normalize.folder_fetch_error(req.url)
    if folder_err:
        job_id = uuid.uuid4().hex
        job = JobRecord(
            job_id=job_id,
            agent_id=req.agent_id,
            pack_id=req.pack_id,
            status=JobStatus.FAILED,
            retention_mode=retention,
            warnings=[folder_err],
        )
        jobs_store.save_job(job)
        return job_to_response(job)

    target_url = url_normalize.to_direct_download_url(req.url)
    try:
        data, filename, content_type = await asyncio.to_thread(
            pack_import.fetch_file_from_url_typed, target_url, _max_upload_bytes()
        )
    except ImportRefError as exc:
        raise HTTPException(status_code=exc.status_code, detail=exc.detail)

    if not data:
        raise HTTPException(status_code=400, detail="Fetched file is empty.")

    # SharePoint/OneDrive anonymous *folder* shares return an HTML browsing page
    # (e.g. onedrive.aspx), not real file bytes. Fail the job with a clear,
    # actionable message instead of a confusing "Unsupported file type".
    html_err = url_normalize.html_response_error(content_type, filename, data)
    if html_err:
        job_id = uuid.uuid4().hex
        job = JobRecord(
            job_id=job_id,
            agent_id=req.agent_id,
            pack_id=req.pack_id,
            status=JobStatus.FAILED,
            retention_mode=retention,
            warnings=[html_err],
        )
        jobs_store.save_job(job)
        return job_to_response(job)

    upload_id = _stage_url_bytes(filename, data)

    ref_req = BuildFromUploadRefRequest(
        upload_id=upload_id,
        agent_id=req.agent_id,
        pack_id=req.pack_id,
        title=req.title,
        source_language=req.source_language,
        canonical_language=req.canonical_language,
        default_classification=req.default_classification,
        retention_mode=req.retention_mode,
        debug_keep_upload=req.debug_keep_upload,
        return_download_url=req.return_download_url,
        delete_source_after_processing=req.delete_source_after_processing,
        translate=req.translate,
    )
    apply_default_classification = "default_classification" in req.model_fields_set

    job_id = uuid.uuid4().hex
    job = JobRecord(
        job_id=job_id,
        agent_id=req.agent_id,
        pack_id=req.pack_id,
        status=JobStatus.PROCESSING,
        retention_mode=retention,
    )
    jobs_store.save_job(job)

    final = await _run_ref_job_bounded(
        job_id, upload_id, ref_req, apply_default_classification, wait_seconds
    )
    if final is not None and final.status in (
        JobStatus.COMPLETED,
        JobStatus.FAILED,
    ):
        # Fast path: build finished within the wait window -> return the full
        # terminal job (download_url included) in a single round-trip.
        return job_to_response(final)
    # Slow path: still running. Return a JobResponse (status PROCESSING) so the
    # response stays a single declared model (download_url is a nullable field
    # in the OpenAPI schema) and callers can poll GET /api/jobs/{job_id}.
    resp = (
        job_to_response(final)
        if final is not None
        else JobResponse(
            job_id=job_id,
            agent_id=req.agent_id,
            pack_id=req.pack_id,
            status=JobStatus.PROCESSING,
            retention_mode=retention,
            message=(
                f"\u23f3 Build in progress (job {job_id}). Call getJob with this "
                f"job_id until status is COMPLETED, then present the download link "
                f"to the user."
            ),
        )
    )
    resp.upload_id = upload_id
    return resp


@router.post(
    "/build-from-upload-ref",
    response_model=BuildRefAcceptedResponse,
    operation_id="buildMemoryPackFromUploadRef",
    include_in_schema=False,
    summary="Build a Safe Memory Pack from a staged upload (async)",
    description=(
        "Build a bilingual pack from a file already staged via /api/uploads "
        "(referenced by upload_id), so LLMs exchange only small JSON. Processing "
        "runs in the background: returns job_id + status immediately; poll GET "
        "/api/jobs/{job_id}. retention_mode sets temporary vs vault storage."
    ),
)
def build_pack_from_upload_ref(
    req: BuildFromUploadRefRequest,
    background: BackgroundTasks,
    auth: UploadAuthContext = Depends(require_upload_or_token),
) -> BuildRefAcceptedResponse:
    """Kick off async pack building from a previously staged upload."""
    # When authorized by a one-time upload link, the claim's server-side
    # settings win over anything the anonymous page put in the request body.
    claim = auth.claim if auth.mode == "token" else None
    # An explicit classification (from the request body, or always present on a
    # one-time upload-link claim) lets us skip per-entry LLM classification.
    apply_default_classification = (
        "default_classification" in req.model_fields_set or claim is not None
    )
    if claim is not None:
        req = req.model_copy(
            update={
                "agent_id": claim.agent_id,
                "pack_id": claim.pack_id,
                "title": claim.title,
                "source_language": claim.source_language,
                "canonical_language": claim.canonical_language,
                "retention_mode": claim.retention_mode,
                "default_classification": _coerce_classification(
                    claim.default_classification
                ),
            }
        )

    record = storage.load_upload_record(req.upload_id)
    if record is None:
        raise HTTPException(status_code=404, detail="Unknown upload_id.")
    if not storage.get_storage().exists(req.upload_id):
        raise HTTPException(
            status_code=409,
            detail="Upload has no content yet. PUT the file bytes first.",
        )
    retention = _coerce_retention_mode(req.retention_mode)

    job_id = uuid.uuid4().hex
    job = JobRecord(
        job_id=job_id,
        agent_id=req.agent_id,
        pack_id=req.pack_id,
        status=JobStatus.PROCESSING,
        retention_mode=retention,
    )
    jobs_store.save_job(job)

    # Bind the job to the claim and consume one use (single-use link).
    if claim is not None:
        upload_links.consume_use(claim, job_id)

    background.add_task(
        _process_ref_job, job_id, req.upload_id, req, apply_default_classification
    )
    return BuildRefAcceptedResponse(
        job_id=job_id, status=JobStatus.PROCESSING.value, upload_id=req.upload_id
    )


def _stage_url_bytes(filename: str, data: bytes) -> str:
    """Stage fetched URL bytes as an upload record and return its upload_id."""
    store = storage.get_storage()
    upload_id = secrets.token_urlsafe(16)
    store.save(upload_id, filename, data)
    record = UploadRecord(
        upload_id=upload_id,
        filename=filename,
        actual_size=len(data),
        status=UploadStatus.RECEIVED,
        upload_token=secrets.token_urlsafe(8),
        expires_at=storage.compute_upload_expires_at(),
        storage_backend=store.backend_name,
        rel_path=store.path_for(upload_id),
    )
    storage.save_upload_record(record)
    return upload_id


@router.post(
    "/build-from-url",
    response_model=JobResponse,
    response_model_exclude_none=False,
    operation_id="buildPackFromUrl",
    summary="Build a Safe Memory Pack from a share link",
    description=(
        "Fetch a file (xlsx, xls, csv, txt, md, docx, pptx, pdf, images) from a "
        "public HTTPS share link (SharePoint, OneDrive, Google Drive, Dropbox) "
        "and build a Safe Memory Pack; scanned PDFs/images are OCR'd. Fast builds "
        "return the finished pack with download_url; slow ones return a job_id "
        "to poll."
    ),
)
async def build_pack_from_url(req: BuildFromUrlRequest):
    """Fetch a file (or folder ZIP) from a URL and build a pack.

    Bounded-synchronous: waits up to ``sync_build_wait_seconds`` for the build to
    finish. If it completes in time the full terminal job (with download_url) is
    returned in one round-trip; otherwise a job_id + PROCESSING is returned for
    the caller to poll GET /api/jobs/{job_id} (backward compatible).
    """
    return await _build_pack_from_url_impl(req)


def _bump_patch(version: str) -> str:
    """Increment the patch component of a semver-like version string."""
    parts = (version or "0.1.0").split(".")
    while len(parts) < 3:
        parts.append("0")
    try:
        parts[2] = str(int(parts[2]) + 1)
    except ValueError:
        parts[2] = "1"
    return ".".join(parts[:3])
